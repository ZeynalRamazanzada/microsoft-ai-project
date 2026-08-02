"""
BDFS — Sunum Oluşturucu (create_presentation.py)
================================================
python-pptx kütüphanesi kullanılarak BDFS projesi için 
otomatik İngilizce PowerPoint sunumu üretir.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Proje kökünü ayarla
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(project_root)

def create_presentation():
    # Yeni bir sunum oluştur
    prs = Presentation()
    
    # Başlık Slaytı (Slide 1)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Behavioral Decision Fatigue Scoring (BDFS)"
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "A Machine Learning Pipeline for Fatigue Detection\n\n"
    subtitle.text += "Authors:\n"
    subtitle.text += "Zeynalabdin Ramazanzade (231805121)\n"
    subtitle.text += "Alperen Sümeroğlu (231805023)\n\n"
    subtitle.text += "Research Question:\nCan temporal and cognitive dynamics accurately predict decision fatigue?"
    
    # Methodology Slaytı (Slide 2)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Methodology — CRISP-DM Pipeline"
    
    tf = body.text_frame
    tf.text = "Business Understanding: Defining decision fatigue metrics."
    
    p = tf.add_paragraph()
    p.text = "Data Understanding & Preparation:"
    p.level = 0
    p2 = tf.add_paragraph()
    p2.text = "Synthetic data generation using Ornstein-Uhlenbeck process."
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "Pre-processing: Imputation, Winsorization, SMOTE."
    p3.level = 1
    
    p4 = tf.add_paragraph()
    p4.text = "Modeling:"
    p4.level = 0
    p5 = tf.add_paragraph()
    p5.text = "Training 5 ML models (LR, RF, XGB, SVM, KNN)."
    p5.level = 1
    p6 = tf.add_paragraph()
    p6.text = "Hyperparameter optimization via Randomized/Grid Search."
    p6.level = 1
    
    p7 = tf.add_paragraph()
    p7.text = "Evaluation & Deployment:"
    p7.level = 0
    p8 = tf.add_paragraph()
    p8.text = "Ablation study, McNemar statistical test, and SHAP interpretability."
    p8.level = 1
    
    # Dataset Summary Slaytı (Slide 3)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Dataset Summary"
    tf = body.text_frame
    tf.text = "Total Records: 150,000 instances (3,000 participants × 50 trials)"
    
    p = tf.add_paragraph()
    p.text = "Total Features: 26"
    p = tf.add_paragraph()
    p.text = "18 raw features (Age, Task Complexity, Reaction Time, DDM parameters, etc.)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "8 engineered features (e.g., rolling_incon_5, fatigue_slope)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Target Variable: fatigue_class (Binary)"
    p = tf.add_paragraph()
    p.text = "Fatigued (1): 33.61%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Non-fatigued (0): 66.39%"
    p.level = 1
    
    # Model Results Table Slaytı (Slide 4)
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Sadece başlık
    title = slide.shapes.title
    title.text = "Model Results Comparison"
    
    # Tablo ekle
    rows, cols = 6, 7
    left, top, width, height = Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.0)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "ROC-AUC", "PR-AUC"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    data = [
        ["LR", "0.886", "0.774", "0.934", "0.847", "0.948", "0.855"],
        ["RF", "0.904", "0.820", "0.916", "0.866", "0.967", "0.923"],
        ["XGB", "0.905", "0.826", "0.909", "0.866", "0.967", "0.921"],
        ["SVM", "0.671", "0.506", "0.865", "0.639", "0.771", "0.531"],
        ["KNN", "0.703", "0.539", "0.811", "0.647", "0.812", "0.658"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx + 1, col_idx).text = cell_data
            
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9.0), Inches(1.0))
    txBox.text_frame.text = "*XGBoost significantly outperformed Logistic Regression (McNemar test, p < 0.001)."
            
    # ROC Curve Slaytı (Slide 5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "ROC Curve Comparison"
    
    img_path = os.path.join("figures", "roc_curves_comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.5), height=Inches(5.5))
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        txBox.text_frame.text = "Image not found: figures/roc_curves_comparison.png"
        
    # Ablation Slaytı (Slide 6)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Ablation Study — Temporal Features are Critical"
    
    img_path = os.path.join("figures", "11_ablation_study.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(2.0), height=Inches(4.5))
    
    txBox = slide.shapes.add_textbox(Inches(6.0), Inches(2.0), Inches(3.5), Inches(4.0))
    tf = txBox.text_frame
    tf.text = "Key Finding:"
    p = tf.add_paragraph()
    p.text = "Removing DDM (cognitive) features barely impacts AUC (-0.0004)."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Removing Temporal features causes a major drop in AUC (-0.0409)."
    p.level = 0
        
    # SHAP Slaytı (Slide 7)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Interpretability: SHAP Feature Ranking"
    
    img_path = os.path.join("figures", "09_shap_summary.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(5.5))
        
    txBox = slide.shapes.add_textbox(Inches(6.5), Inches(2.0), Inches(3.0), Inches(4.0))
    tf = txBox.text_frame
    tf.text = "Top Features:"
    p = tf.add_paragraph()
    p.text = "1. rolling_incon_5"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "2. rolling_incon_10"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "3. pref_reversal_rate"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Our engineered temporal features dominate the model's decision process."
    p.level = 0
    
    # Conclusion Slaytı (Slide 8)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Conclusion & Academic Contribution"
    
    tf = body.text_frame
    tf.text = "Machine learning effectively captures behavioral decision fatigue."
    
    p = tf.add_paragraph()
    p.text = "XGBoost and Random Forest achieve near-perfect discrimination (AUC ~0.967) even in imbalanced settings."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Paradigm Shift in Fatigue Analysis:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Contrary to traditional psychological models that rely solely on static cognitive parameters (DDM), this study proves that dynamic, temporal trends (e.g., rolling inconsistencies) are significantly stronger predictors."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Future Work:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Real-time deployment of this pipeline for occupational fatigue monitoring."
    p.level = 1
    
    # Kaydet
    output_path = "BDFS_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
